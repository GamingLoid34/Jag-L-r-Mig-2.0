import streamlit as st
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation
from gtts import gTTS
import tempfile
import os
import json
import time # För att tidstämpla framsteg

# --- INSTÄLLNINGAR ---
st.set_page_config(page_title="Jag Lär Mig", page_icon="📖", layout="wide")
# --- SÄKER START AV MINNET ---
if "subjects" not in st.session_state:
    st.session_state.subjects = {"Allmänt": {"material": "", "history": []}}
if "current_subject" not in st.session_state:
    st.session_state.current_subject = "Allmänt"
if "flashcards" not in st.session_state:
    st.session_state.flashcards = {}
# -----------------------------
# --- MAPPING: ÄMNEN TILL BAKGRUNDSBILDER (som du ville ha) ---
BACKGROUND_MAP = {
    "NO": "url('https://images.unsplash.com/photo-1582719478253-6ce7ebdf11c8?q=80&w=2500&auto=format&fit=crop&ixlib=rb-4.0.3&ixid=M3wxMjA3fDB8MHxwaG90by1wYWdlfHx8fGVufDB8fHx8fA%3D%3D')",
    "Geografi": "url('https://images.unsplash.com/photo-1541334311090-344070a7b055?q=80&w=2500&auto=format&fit=crop&ixlib=rb-4.0.3&ixid=M3wxMjA3fDB8MHxwaG90by1wYWdlfHx8fGVufDB8fHx8fA%3D%3D')",
    "Idrott": "url('https://images.unsplash.com/photo-1517590806450-482a2a...')", # Kortad länk
    "Matte": "url('https://images.unsplash.com/photo-1596495689108-bc31c626456f?q=80&w=2500&auto=format&fit=crop&ixlib=rb-4.0.3&ixid=M3wxMjA3fDB8MHxwaG90by1wYWdlfHx8fGVufDB8fHx8fA%3D%3D')",
    "Allmänt": "url('https://images.unsplash.com/photo-1507525428034-b723cf961d3e?q=80&w=2500&auto=format&fit=crop&ixlib=rb-4.0.3&ixid=M3wxMjA3fDB8MHxwaG90by1wYWdlfHx8fGVufDB8fHx8fA%3D%3D')",
}

def set_background(subject_name):
    # Hämtar URL från mappen
    bg_url = BACKGROUND_MAP.get(subject_name, BACKGROUND_MAP['Allmänt'])
    
    # Injekterar CSS för bakgrundsbild
    st.markdown(
        f"""
        <style>
        .stApp {{
            background-image: {bg_url};
            background-size: cover;
            background-repeat: no-repeat;
            background-attachment: fixed;
            transition: background-image 0.5s ease-in-out;
            color: white; 
        }}
        </style>
        """,
        unsafe_allow_html=True
    )

# --- TRACKING: NY DATASTRUKTUR FÖR FRAMSTEG ---
def initialize_tracking():
    if "subjects" not in st.session_state:
        st.session_state.subjects = {"Allmänt": {"material": "", "history": []}} # NYTT: history lista
    
    # Skapa tracking-struktur för alla ämnen
    for sub in st.session_state.subjects:
        if isinstance(st.session_state.subjects[sub], str): # Fixa gammal struktur
            st.session_state.subjects[sub] = {"material": st.session_state.subjects[sub], "history": []}
        if "history" not in st.session_state.subjects[sub]:
            st.session_state.subjects[sub]["history"] = []

# Kör initialisering av tracking
initialize_tracking()
# ---------------------------------------------


# --- FUNKTIONER ---

# (Dessa är samma som i din tidigare kod)
def extract_text_from_pdf(pdf_file):
    text = ""
    reader = PdfReader(pdf_file)
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def generate_speech_simple(text):
    try:
        tts = gTTS(text=text, lang='sv')
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
            tts.save(fp.name)
            return fp.name
    except Exception as e:
        st.error(f"Kunde inte skapa ljud: {e}")
        return None

def get_gemini_response(prompt, context, api_key):
    if not api_key: return "⚠️ Fel: API-nyckel saknas. Lägg in nyckeln i Secrets!"
    
    # Hämtar nyckel från secrets nu när vi har bytt till den robusta lösningen
    gemini_key = st.secrets.get("GEMINI_API_KEY", api_key)
    if not gemini_key: return "⚠️ Fel: API-nyckel saknas!"

    genai.configure(api_key=gemini_key) 
    
    system_instruction = (
        "Du är en smart och pedagogisk studiecoach i appen 'Jag Lär Mig'. "
        "Skapa alltid output som en ren JSON-lista om det begärs. "
        "Din uppgift är att hjälpa användaren att förstå sitt studiematerial."
    )
    model = genai.GenerativeModel('gemini-1.5-flash', system_instruction=system_instruction)
    
    full_prompt = f"Studiematerial:\n{context}\n\nUppgift/Fråga: {prompt}"
    
    try:
        return model.generate_content(full_prompt).text
    except Exception as e:
        st.error(f"❌ AI-anropet misslyckades. Nyckelproblem eller modellkvot.")
        return f"Fel: {e}"


# --- SIDOPANEL (MENY) ---

# Global api_key variabel (hämtas från Secrets)
api_key = st.secrets.get("GEMINI_API_KEY", None)
if not api_key:
    # Om nyckeln saknas i Secrets, visa varningsruta
    with st.sidebar:
        st.warning("⚠️ Varning: API-nyckel saknas. Lägg in den i Streamlit Secrets!")
        
with st.sidebar:
    st.title("📖 Jag Lär Mig")
    with st.sidebar:
    st.title("📖 Jag Lär Mig")
    
    # --- DIAGNOS KOD (TA BORT SEN) ---
    if "GEMINI_API_KEY" in st.secrets:
        nyckel = st.secrets["GEMINI_API_KEY"]
        st.info(f"🔍 Debug: Appen ser nyckeln: '{nyckel[:5]}...' (Längd: {len(nyckel)})")
    else:
        st.error("Debug: Appen hittar INGEN nyckel i Secrets!")
    # ---------------------------------

    # ... resten av koden ...
    if api_key:
        st.success("🔑 Nyckel laddad från Secrets!")

    st.divider()

    st.subheader("📂 Mina Ämnen")
    subject_list = list(st.session_state.subjects.keys())
    
    # Välj ämne
    selected_sub = st.selectbox("Välj ämne att plugga:", subject_list, index=subject_list.index(st.session_state.current_subject))
    st.session_state.current_subject = selected_sub
    
    # Skapa nytt ämne
    new_sub = st.text_input("Lägg till nytt ämne:")
    if st.button("Skapa mapp") and new_sub:
        st.session_state.subjects[new_sub] = {"material": "", "history": []}
        st.session_state.current_subject = new_sub
        st.success(f"Mappen '{new_sub}' skapad!")
        st.rerun()

    st.divider()
    
    # Uppladdning
    st.subheader(f"📥 Ladda upp till: {st.session_state.current_subject}")
    uploaded_files = st.file_uploader("Släpp filer här (PDF, PPTX)", accept_multiple_files=True)
    
    if st.button("Spara materialet"):
        current_data = st.session_state.subjects[st.session_state.current_subject]["material"]
        count = 0
        for file in uploaded_files:
            if file.name.endswith(".pdf"):
                current_data += f"\n--- {file.name} ---\n" + extract_text_from_pdf(file)
                count += 1
            elif file.name.endswith(".pptx"):
                current_data += f"\n--- {file.name} ---\n" + extract_text_from_pptx(file)
                count += 1
        
        st.session_state.subjects[st.session_state.current_subject]["material"] = current_data
        st.success(f"Sparade {count} filer i {st.session_state.current_subject}!")


# --- HUVUDVY ---
# Ladda bakgrund (Baseras på valt ämne)
set_background(st.session_state.current_subject) 

st.header(f"Studerar: {st.session_state.current_subject}")

current_data = st.session_state.subjects[st.session_state.current_subject]
current_material = current_data["material"]
current_history = current_data["history"]


if not current_material:
    st.info("👈 Den här mappen är tom. Börja med att ladda upp material i menyn!")
else:
    # --- FLIKAR ---
    tab1, tab2, tab3, tab4 = st.tabs(["📝 Material", "🧠 Förhör", "📊 Framsteg & Flashcards", "🎧 Lyssna"])


    # FLIK 1: REDIGERA & STRUKTUR
    with tab1:
        st.subheader("Ditt material")
        edited_text = st.text_area("Innehåll", current_material, height=300)
        
        if st.button("Spara ändringar i texten"):
            st.session_state.subjects[st.session_state.current_subject]["material"] = edited_text
            st.success("Uppdaterat!")
            st.rerun()

        st.divider()
        if st.button("✨ Dela upp texten i kapitel (AI)"):
            with st.spinner("Analyserar struktur..."):
                chapters = get_gemini_response(
                    "Dela upp texten i tydliga kapitel/avsnitt med rubriker.", 
                    edited_text, api_key
                )
                st.markdown(chapters)

    # FLIK 2: FÖRHÖR (Standard chatt)
    with tab2:
        st.subheader("Plugga med AI")
        
        c1, c2 = st.columns(2)
        if c1.button("Skapa prov"):
            with st.spinner("Skapar prov..."):
                test = get_gemini_response("Skapa ett prov med 5 frågor + facit.", current_material, api_key)
                st.markdown(test)
                # OBS: Här lägger vi till logik för att spara resultatet när användaren klickar 'Rätta'

        if c2.button("Sammanfatta"):
            with st.spinner("Sammanfattar..."):
                summary = get_gemini_response("Sammanfatta det viktigaste i punktform.", current_material, api_key)
                st.markdown(summary)

        st.divider()
        user_q = st.chat_input("Ställ en fråga om materialet...")
        if user_q:
            st.chat_message("user").write(user_q)
            with st.spinner("Tänker..."):
                ans = get_gemini_response(user_q, current_material, api_key)
                st.chat_message("assistant").write(ans)
    
    # ----------------------------------------------------
    # NY FLIK: TRACKING & FLASHCARDS (Kärnan i Adaptivitet)
    # ----------------------------------------------------
    with tab3:
        st.header("Framsteg & Interaktion")
        
        # 1. FRAMSTEGSÖVERSIKT (Tracking)
        st.subheader("1. Översikt")
        st.info(f"Hittills sparade resultat för {st.session_state.current_subject}: {len(current_history)} sessioner.")
        if current_history:
            st.write(current_history[-1]) # Visa senaste resultatet
        
        st.divider()

        # 2. FLASHCARDS (Interaktivt läge)
        st.subheader("2. Flashcards (Spaced Repetition)")
        
        # Initiera state för Flashcards
        if "flashcards" not in st.session_state or st.session_state.current_subject not in st.session_state.flashcards:
             st.session_state.flashcards = {st.session_state.current_subject: None}

        # A) KNAPP FÖR GENERERING
        if st.button("▶️ Generera 5 nya Flashcards"):
            with st.spinner("AI:n skapar kort..."):
                # Be AI:n skapa JSON-format för enkel hantering
                cards_json = get_gemini_response(
                    "Skapa 5 flashcards (fråga/svar-par) från materialet. Svara ENDAST med en ren JSON-lista i formatet: [{'question': '...', 'answer': '...'}, ...]",
                    current_material, api_key
                )
                # Försök parsa JSON outputen från AI:n
                try:
                    cards = json.loads(cards_json)
                    st.session_state.flashcards[st.session_state.current_subject] = cards
                    st.success("5 kort skapade! Börja öva nedan.")
                except json.JSONDecodeError:
                    st.error("AI:n gav fel format. Försök igen eller justera prompten.")
                    st.write(cards_json) # Visa outputen för debugging

        # B) FLASHCARD UI
        cards = st.session_state.flashcards.get(st.session_state.current_subject)
        
        if cards:
            st.write(f"Du har {len(cards)} kort att öva på.")
            card_col, score_col = st.columns([3, 1])

            # Visa ett kort i taget
            if "card_index" not in st.session_state:
                st.session_state.card_index = 0
            
            card_index = st.session_state.card_index
            if card_index < len(cards):
                card = cards[card_index]

                with card_col:
                    st.markdown(f"### Kort {card_index + 1} av {len(cards)}")
                    st.info(card['question'])
                    
                    if st.button("Visa svar"):
                        st.success(card['answer'])

                    # Markera resultat och gå vidare
                    col_know, col_forget, _ = st.columns(3)
                    if col_know.button("Kunde den (👍)"):
                        # Lägg till logik för Spaced Repetition/Tracking här
                        st.session_state.card_index += 1
                        st.rerun()
                    if col_forget.button("Glömde den (👎)"):
                        # Lägg till logik för att lägga kortet sist i repetitionen här
                        st.session_state.card_index += 1
                        st.rerun()
                
            else:
                st.success("Bra jobbat! Du har gått igenom alla kort för denna gång.")
                if st.button("Börja om"):
                     st.session_state.card_index = 0
                     st.rerun()


    # FLIK 4: LYSSNA
    with tab4:
        st.subheader("Uppläsning")
        
        text_to_read = st.text_area("Text att läsa upp:", value=current_material[:3000], height=150)

        if st.button("▶️ Spela upp"):
            with st.spinner("Skapar ljud..."):
                audio_path = generate_speech_simple(text_to_read)
                if audio_path:
                    st.audio(audio_path, format="audio/mp3")
